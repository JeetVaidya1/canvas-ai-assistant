# ingest.py — GPT-5 optimized ingestion (cleaning, chunking, embeddings, metadata, upsert)
import os, io, re, json, time, hashlib, tempfile
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from dotenv import load_dotenv
from openai import OpenAI
from vector_store import VectorStore

# Optional imports (fail gracefully if unavailable)
try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from docx import Document
except Exception:
    Document = None

# OCR (optional)
try:
    import pytesseract
    from PIL import Image
except Exception:
    pytesseract, Image = None, None

# ── Env / clients ────────────────────────────────────────────────────────────
load_dotenv()
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
vector_store = VectorStore()

# Supabase (used by delete helpers)
from supabase import create_client
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

# ── Config knobs (override via .env) ─────────────────────────────────────────
EMBED_MODEL = os.getenv("EMBEDDINGS_MODEL", "text-embedding-3-large")
CHUNK_CHARS = int(os.getenv("CHUNK_CHARS", "1800"))        # ~1k-1.4k tokens depending on language
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "250"))
EMBED_BATCH = int(os.getenv("EMBED_BATCH", "64"))
ENABLE_OCR = os.getenv("ENABLE_OCR", "false").lower() == "true"
SUPPORTED = {".pdf", ".pptx", ".docx", ".txt", ".md", ".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}

# ── Data classes ─────────────────────────────────────────────────────────────
@dataclass
class RawBlock:
    text: str
    page: Optional[int] = None
    slide: Optional[int] = None
    section: Optional[str] = None

# ── Utils ────────────────────────────────────────────────────────────────────
def _norm_ws(s: str) -> str:
    s = s.replace("\x00", " ")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def _dedup_blocks(blocks: List[RawBlock]) -> List[RawBlock]:
    seen = set()
    out = []
    for b in blocks:
        t = b.text.strip()
        if not t:
            continue
        h = hashlib.sha256(t.encode("utf-8")).hexdigest()
        if h in seen:
            continue
        seen.add(h)
        out.append(b)
    return out

def _chunk_text(text: str, chunk_chars: int, overlap: int) -> List[str]:
    text = text.strip()
    if len(text) <= chunk_chars:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_chars)
        # Try to break on a sentence boundary
        slice_ = text[start:end]
        m = re.search(r"(?s)^(.+?)(?:\n\n|(?<=[.!?])\s)", slice_)
        if m and len(m.group(0)) > chunk_chars * 0.6:
            slice_ = m.group(0)
        chunks.append(slice_.strip())
        if end == len(text):
            break
        start += max(1, len(slice_) - overlap)
    return [c for c in chunks if len(c) > 20]

def _to_chunks(blocks: List[RawBlock]) -> List[Dict[str, Any]]:
    """Chunk per block; carry page/slide/section metadata into each chunk."""
    out = []
    chunk_id = 0
    for b in blocks:
        parts = _chunk_text(b.text, CHUNK_CHARS, CHUNK_OVERLAP)
        for p in parts:
            chunk_id += 1
            out.append({
                "content": p,
                "page": b.page,
                "slide": b.slide,
                "section": b.section,
                "chunk_id": chunk_id
            })
    return out

def _embed_batch(texts: List[str], retry=3, backoff=1.5) -> List[List[float]]:
    out = []
    for i in range(0, len(texts), EMBED_BATCH):
        batch = texts[i:i+EMBED_BATCH]
        for attempt in range(retry):
            try:
                resp = openai_client.embeddings.create(model=EMBED_MODEL, input=batch)
                out.extend([d.embedding for d in resp.data])
                break
            except Exception as e:
                if attempt == retry - 1:
                    raise
                time.sleep(backoff ** (attempt + 1))
    return out

# ── Extractors ───────────────────────────────────────────────────────────────
def extract_pdf(file_bytes: bytes) -> List[RawBlock]:
    if not pdfplumber:
        return []
    blocks: List[RawBlock] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            try:
                txt = page.extract_text() or ""
            except Exception:
                txt = ""
            txt = _norm_ws(txt)
            if txt:
                blocks.append(RawBlock(text=txt, page=i, section=f"Page {i}"))
    return _dedup_blocks(blocks)

def extract_pptx(file_bytes: bytes) -> List[RawBlock]:
    if not Presentation:
        return []
    prs = Presentation(io.BytesIO(file_bytes))
    blocks: List[RawBlock] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            # text boxes
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            # tables
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_txt = " | ".join([cell.text for cell in row.cells if cell.text])
                    if row_txt.strip():
                        texts.append(row_txt)
            # notes
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                texts.append(f"(Speaker Notes) {notes}")
        merged = _norm_ws("\n".join(t for t in texts if t))
        if merged:
            blocks.append(RawBlock(text=merged, slide=i, section=f"Slide {i}"))
    return _dedup_blocks(blocks)

def extract_docx(file_bytes: bytes) -> List[RawBlock]:
    if not Document:
        return []
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(file_bytes)
        path = tmp.name
    try:
        doc = Document(path)
        cur_section = None
        paras: List[RawBlock] = []
        for p in doc.paragraphs:
            txt = _norm_ws(p.text or "")
            if not txt:
                continue
            style = (p.style.name if p.style else "") or ""
            if "Heading" in style:
                cur_section = txt
                continue
            paras.append(RawBlock(text=txt, section=cur_section))
        # Merge small paragraphs under a section to form bigger blocks
        blocks: List[RawBlock] = []
        buf = []
        cur = None
        for rb in paras:
            if rb.section != cur and buf:
                blocks.append(RawBlock(text=_norm_ws("\n".join(buf)), section=cur))
                buf = []
            cur = rb.section
            buf.append(rb.text)
        if buf:
            blocks.append(RawBlock(text=_norm_ws("\n".join(buf)), section=cur))
        return _dedup_blocks(blocks)
    finally:
        try: os.remove(path)
        except Exception: pass

def extract_txt_like(file_bytes: bytes) -> List[RawBlock]:
    try:
        txt = file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        txt = ""
    txt = _norm_ws(txt)
    if not txt:
        return []
    # Split by big gaps/headers to create sections
    parts = re.split(r"\n\s*\n", txt)
    blocks = [RawBlock(text=_norm_ws(p)) for p in parts if p.strip()]
    return _dedup_blocks(blocks)

def extract_image_ocr(file_bytes: bytes) -> List[RawBlock]:
    if not (ENABLE_OCR and pytesseract and Image):
        return []
    try:
        img = Image.open(io.BytesIO(file_bytes))
        txt = pytesseract.image_to_string(img)
        txt = _norm_ws(txt)
        return [RawBlock(text=txt, section="Image OCR")] if txt else []
    except Exception:
        return []

# ── Dispatcher ───────────────────────────────────────────────────────────────
def _extract_by_ext(filename: str, file_bytes: bytes) -> List[RawBlock]:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return extract_pdf(file_bytes)
    if ext == ".pptx":
        return extract_pptx(file_bytes)
    if ext == ".docx":
        return extract_docx(file_bytes)
    if ext in {".txt", ".md"}:
        return extract_txt_like(file_bytes)
    if ext in {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}:
        return extract_image_ocr(file_bytes)
    return []

# ── Public API ────────────────────────────────────────────────────────────────
def process_file(filename: str, file_bytes: bytes, course_id: str) -> List[Dict[str, Any]]:
    """
    Ingest a file: extract text, chunk, embed in batches, and store in the vector DB.
    Returns a preview list of chunks.
    """
    ext = os.path.splitext(filename)[1].lower()
    if ext not in SUPPORTED:
        return [{"chunk": f"Unsupported file type: {ext}"}]

    # 1) Extract blocks
    blocks = _extract_by_ext(filename, file_bytes)
    if not blocks:
        return [{"chunk": f"No usable text found in {filename}"}]

    # 2) Chunk with metadata
    chunk_dicts = _to_chunks(blocks)

    # 3) Clean & dedup short/empty chunks
    clean_chunks = []
    seen_hash = set()
    for c in chunk_dicts:
        content = _norm_ws(c["content"])
        if len(content) < 40:
            continue
        h = hashlib.sha256(content.encode("utf-8")).hexdigest()
        if h in seen_hash:
            continue
        seen_hash.add(h)
        c["content"] = content
        c["sha256"] = h
        clean_chunks.append(c)

    if not clean_chunks:
        return [{"chunk": f"No usable text after cleaning in {filename}"}]

    # 4) Embed in batches
    texts = [c["content"] for c in clean_chunks]
    embeddings = _embed_batch(texts)

    # 5) Store each chunk + vector (with rich metadata)
    for c, emb in zip(clean_chunks, embeddings):
        vector_store.add(
            course_id=course_id,
            doc_name=filename,
            chunk_id=c["chunk_id"],
            embedding=emb,
            content=c["content"],
            page=c.get("page"),
            slide=c.get("slide"),
            section=c.get("section"),
            sha256=c.get("sha256")
        )

    # 6) (Optional) store file metadata summary for UI (first/last page, counts)
    try:
        supabase.table("files").upsert({
            "course_id": course_id,
            "filename": filename,
            "ext": ext,
            "num_chunks": len(clean_chunks)
        }, on_conflict="course_id,filename").execute()
    except Exception as e:
        print(f"files upsert warn: {e}")

    # 7) Return preview of first 5 chunks w/ metadata
    preview = []
    for c in clean_chunks[:5]:
        meta = []
        if c.get("page"): meta.append(f"p{c['page']}")
        if c.get("slide"): meta.append(f"slide{c['slide']}")
        if c.get("section"): meta.append(c["section"])
        preview.append({
            "chunk": c["content"],
            "meta": " • ".join(meta) if meta else ""
        })
    return preview

# ── Delete helpers (kept, with small safety tweaks) ──────────────────────────
def delete_file_from_course(course_id: str, filename: str) -> bool:
    """Delete all embeddings + storage + metadata for a specific file from a course."""
    try:
        supabase.table("embeddings").delete().eq("course_id", course_id).eq("doc_name", filename).execute()
        supabase.table("files").delete().eq("course_id", course_id).eq("filename", filename).execute()
        # Storage
        try:
            supabase.storage.from_("course-files").remove([f"{course_id}/{filename}"])
        except Exception as e:
            print(f"Storage deletion warn: {e}")
        return True
    except Exception as e:
        print(f"Error deleting file {filename} from course {course_id}: {e}")
        return False

def delete_course(course_id: str) -> bool:
    """Delete an entire course and all its data."""
    try:
        supabase.table("embeddings").delete().eq("course_id", course_id).execute()
        supabase.table("files").delete().eq("course_id", course_id).execute()

        sessions_resp = supabase.table("chat_sessions").select("id").eq("course_id", course_id).execute()
        session_ids = [s["id"] for s in (sessions_resp.data or [])]
        for session_id in session_ids:
            supabase.table("messages").delete().eq("session_id", session_id).execute()
        supabase.table("chat_sessions").delete().eq("course_id", course_id).execute()

        supabase.table("courses").delete().eq("course_id", course_id).execute()

        # Storage cleanup
        try:
            files_list = supabase.storage.from_("course-files").list(course_id) or []
            if files_list:
                supabase.storage.from_("course-files").remove([f"{course_id}/{f['name']}" for f in files_list])
        except Exception as e:
            print(f"Storage cleanup warn: {e}")

        return True
    except Exception as e:
        print(f"Error deleting course {course_id}: {e}")
        return False
