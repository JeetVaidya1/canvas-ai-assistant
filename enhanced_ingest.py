# enhanced_ingest.py — GPT-5 optimized multimodal ingest
import os, io, tempfile, base64, hashlib, re
from typing import List, Dict, Any, Optional
import fitz  # PyMuPDF
import pdfplumber
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from openai import OpenAI
from vector_store import VectorStore

load_dotenv()
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
vector_store = VectorStore()

EMBED_MODEL = os.getenv("EMBEDDINGS_MODEL", "text-embedding-3-small")
CHUNK_CHARS = int(os.getenv("CHUNK_CHARS", "1800"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "250"))
MAX_IMAGES = int(os.getenv("INGEST_MAX_IMAGES", "5"))

def _norm_ws(s: str) -> str:
    s = s.replace("\x00", " ")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def _chunk(text: str) -> List[str]:
    text = (text or "").strip()
    if not text:
        return []
    if len(text) <= CHUNK_CHARS:
        return [text]
    chunks, start = [], 0
    while start < len(text):
        end = min(len(text), start + CHUNK_CHARS)
        slice_ = text[start:end]
        # try not to break mid-sentence
        cut = max(slice_.rfind("\n\n"), slice_.rfind(". "), slice_.rfind("? "), slice_.rfind("! "))
        if cut > CHUNK_CHARS * 0.4:
            slice_ = slice_[:cut+1]
        chunks.append(slice_.strip())
        start += max(1, len(slice_) - CHUNK_OVERLAP)
    return [c for c in chunks if len(c) > 40]

def extract_text_from_pdf_basic(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Return per-page blocks [{'text':..., 'page': int}] for metadata."""
    blocks = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            txt = _norm_ws(page.extract_text() or "")
            if txt:
                blocks.append({"text": txt, "page": i, "section": f"Page {i}"})
    return blocks

def extract_images_from_pdf(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Extract images (png base64) with page index."""
    out = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for pidx in range(len(doc)):
            page = doc[pidx]
            for j, img in enumerate(page.get_images(full=True)):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n - pix.alpha < 4:  # grayscale/rgb
                        data = pix.tobytes("png")
                        out.append({
                            "b64": base64.b64encode(data).decode(),
                            "page": pidx + 1,
                            "index": j
                        })
                except Exception as e:
                    print(f"image extract warn p{pidx+1} #{j}: {e}")
        doc.close()
    except Exception as e:
        print(f"pdf image pass warn: {e}")
    return out

def extract_text_from_pptx(file_bytes: bytes) -> List[Dict[str, Any]]:
    prs = Presentation(io.BytesIO(file_bytes))
    blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_txt = " | ".join([cell.text for cell in row.cells if cell.text])
                    if row_txt.strip(): texts.append(row_txt)
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes: texts.append(f"(Speaker Notes) {notes}")
        merged = _norm_ws("\n".join(t for t in texts if t))
        if merged:
            blocks.append({"text": merged, "slide": i, "section": f"Slide {i}"})
    return blocks

def extract_text_from_docx(file_bytes: bytes) -> List[Dict[str, Any]]:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(file_bytes)
        path = tmp.name
    try:
        doc = Document(path)
        blocks, buf, cur = [], [], None
        for p in doc.paragraphs:
            t = _norm_ws(p.text or "")
            if not t: continue
            style = (p.style.name if p.style else "") or ""
            if "Heading" in style:
                # flush previous section
                if buf:
                    blocks.append({"text": _norm_ws("\n".join(buf)), "section": cur})
                    buf = []
                cur = t
            else:
                buf.append(t)
        if buf:
            blocks.append({"text": _norm_ws("\n".join(buf)), "section": cur})
        return blocks
    finally:
        try: os.remove(path)
        except: pass

def describe_image_with_gpt(image_b64: str, context: str = "") -> str:
    """Use GPT-5 (vision) for a crisp educational description."""
    try:
        prompt = (
            "Describe this image for a student: what it shows, the key concepts, "
            "and any labels/axes. Be concise but specific. " + (f"\nContext: {context}" if context else "")
        )
        resp = openai_client.chat.completions.create(
            model=os.getenv("MODEL_COMPLEX", "gpt-5"),
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}" }}
                ]
            }],
            temperature=0.2,
            max_tokens=400
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        print(f"vision describe warn: {e}")
        return "Image description unavailable."

def _sha(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8")).hexdigest()

def process_file_enhanced(filename: str, file_bytes: bytes, course_id: str) -> List[Dict[str, Any]]:
    """Enhanced processing for PDFs/PPTX/DOCX with text + (limited) image pass."""
    ext = os.path.splitext(filename)[1].lower()
    text_blocks: List[Dict[str, Any]] = []
    image_infos: List[Dict[str, Any]] = []

    if ext == ".pdf":
        text_blocks = extract_text_from_pdf_basic(file_bytes)
        image_infos = extract_images_from_pdf(file_bytes)
    elif ext == ".pptx":
        text_blocks = extract_text_from_pptx(file_bytes)
        image_infos = []  # PPTX image OCR optional later
    elif ext == ".docx":
        text_blocks = extract_text_from_docx(file_bytes)
        image_infos = []
    else:
        return [{"chunk": f"Unsupported file type: {ext}"}]

    # 1) Text chunks → embeddings
    preview = []
    chunk_id = 0
    for blk in text_blocks:
        for piece in _chunk(blk.get("text", "")):
            chunk_id += 1
            sha = _sha(piece)
            # embed
            emb = openai_client.embeddings.create(model=EMBED_MODEL, input=[piece]).data[0].embedding
            vector_store.add(
                course_id=course_id,
                doc_name=filename,
                chunk_id=chunk_id,
                embedding=emb,
                content=piece,
                page=blk.get("page"),
                slide=blk.get("slide"),
                section=blk.get("section"),
                sha256=sha
            )
            if len(preview) < 3:
                meta_bits = []
                if blk.get("page"): meta_bits.append(f"p{blk['page']}")
                if blk.get("slide"): meta_bits.append(f"slide{blk['slide']}")
                if blk.get("section"): meta_bits.append(str(blk["section"]))
                preview.append({"type": "text", "chunk": piece[:200] + ("..." if len(piece) > 200 else ""), "meta": " • ".join(meta_bits)})

    # 2) Limited images → describe → embed description
    if image_infos:
        for i, info in enumerate(image_infos[:MAX_IMAGES]):
            try:
                desc = describe_image_with_gpt(info["b64"], context=f"{filename} page {info.get('page')}")
                text_blob = f"[IMAGE CONTENT — Page {info.get('page')}] {desc}"
                chunk_id += 1
                sha = _sha(text_blob)
                emb = openai_client.embeddings.create(model=EMBED_MODEL, input=[text_blob]).data[0].embedding
                vector_store.add(
                    course_id=course_id,
                    doc_name=filename,
                    chunk_id=chunk_id,
                    embedding=emb,
                    content=text_blob,
                    page=info.get("page"),
                    slide=None,
                    section="Image",
                    sha256=sha
                )
                if len(preview) < 5:
                    preview.append({"type": "image", "chunk": desc[:200] + ("..." if len(desc) > 200 else ""), "meta": f"p{info.get('page')}"})
            except Exception as e:
                print(f"image ingest warn: {e}")

    return preview
